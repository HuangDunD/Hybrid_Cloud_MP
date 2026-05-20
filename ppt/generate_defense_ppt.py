#!/usr/bin/env python3
from __future__ import annotations

import re
import zipfile
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PAPER = ROOT / "docs" / "支持多写OLTP数据库一体机页面亲和性论文_修改版.docx"
TEMPLATE = ROOT / "ppt" / "一抹蓝.pptx"
OUT = ROOT / "ppt" / "支持多写OLTP数据库一体机页面亲和性答辩.pptx"
ASSET_DIR = ROOT / "ppt" / "generated_assets"

FONT = "Microsoft YaHei"
BLUE = RGBColor(23, 87, 151)
DARK = RGBColor(16, 43, 78)
TEXT = RGBColor(35, 46, 61)
MUTED = RGBColor(91, 108, 128)
LIGHT = RGBColor(236, 245, 253)
GRID = RGBColor(203, 221, 239)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(19, 147, 96)
ORANGE = RGBColor(224, 123, 35)
RED = RGBColor(196, 60, 60)
GRAY = RGBColor(248, 251, 254)
SHADE = RGBColor(225, 237, 250)


def emu(v: float) -> int:
    return Inches(v)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = list(slide_id_list)[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=WHITE, line=GRID, radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(slide, content: str, x, y, w, h, size=13, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = content
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return shape


def bullets(slide, lines, x, y, w, h, size=11.5, color=TEXT, spacing=4):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(spacing)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return shape


def header(slide, title: str, section: str) -> None:
    rect(slide, emu(0), emu(0), emu(13.333), emu(0.62), fill=BLUE, line=BLUE)
    text(slide, title, emu(0.48), emu(0.15), emu(10.7), emu(0.32),
         size=19, color=WHITE, bold=True)
    text(slide, section, emu(11.25), emu(0.19), emu(1.55), emu(0.2),
         size=10.5, color=WHITE, align=PP_ALIGN.RIGHT)
    rect(slide, emu(0), emu(7.27), emu(13.333), emu(0.23), fill=LIGHT, line=LIGHT)
    text(slide, "支持多写 OLTP 数据库一体机页面亲和性关键技术研究",
         emu(0.5), emu(7.31), emu(5.8), emu(0.12), size=7.2, color=MUTED)
    text(slide, "CDUT · WookongDB-MP",
         emu(10.3), emu(7.31), emu(2.5), emu(0.12), size=7.2, color=MUTED, align=PP_ALIGN.RIGHT)


def lead(slide, content: str) -> None:
    """Evaluator-facing one-liner directly under the header.
    Explains in plain language what the slide is about, for reviewers who are
    not specialists in shared-storage multi-write databases.
    """
    text(slide, content, emu(0.55), emu(0.68), emu(12.23), emu(0.26),
         size=10.2, color=MUTED, italic=True)


def conclusion(slide, content: str, bridge: str = "") -> None:
    """Bottom bar: GREEN bold takeaway + (optional) MUTED italic bridge phrase."""
    h = 0.66 if bridge else 0.46
    rect(slide, emu(0.55), emu(6.42), emu(12.23), emu(h), fill=LIGHT, line=GRID)
    text(slide, content, emu(0.78), emu(6.50), emu(11.77), emu(0.22),
         size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    if bridge:
        text(slide, "▸ " + bridge, emu(0.78), emu(6.80), emu(11.77), emu(0.20),
             size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


def image_fit(slide, path, x, y, w, h, caption=None, pad=0.08):
    rect(slide, x, y, w, h, fill=WHITE, line=GRID)
    if not path or not Path(path).exists():
        text(slide, "图片缺失", x, y + h // 2, w, emu(0.25),
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    cap_h = emu(0.22) if caption else 0
    max_w = w - emu(pad * 2)
    max_h = h - emu(pad * 2) - cap_h
    scale = min(max_w / iw, max_h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px = int(x + (w - pw) / 2)
    py = int(y + emu(pad) + (max_h - ph) / 2)
    slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)
    if caption:
        text(slide, caption, x + emu(0.1), y + h - emu(0.25), w - emu(0.2), emu(0.14),
             size=7.5, color=MUTED, align=PP_ALIGN.CENTER)


def side_points(slide, title, lines, x, y, w, h, size=11.4):
    rect(slide, x, y, w, h, fill=GRAY, line=GRID)
    text(slide, title, x + emu(0.22), y + emu(0.22), w - emu(0.44), emu(0.28),
         size=15, color=DARK, bold=True)
    bullets(slide, lines, x + emu(0.25), y + emu(0.72), w - emu(0.5), h - emu(0.9), size=size)


def metric(slide, value, label, x, y, w, h, color=BLUE):
    rect(slide, x, y, w, h, fill=WHITE, line=GRID)
    text(slide, value, x + emu(0.08), y + emu(0.12), w - emu(0.16), emu(0.3),
         size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    text(slide, label, x + emu(0.08), y + emu(0.5), w - emu(0.16), emu(0.2),
         size=8.8, color=MUTED, align=PP_ALIGN.CENTER)


def table(slide, rows, x, y, w, h, widths=None, size=9.6, highlight_col=None):
    cols = len(rows[0])
    widths = widths or [1] * cols
    total = sum(widths)
    col_ws = [int(w * weight / total) for weight in widths]
    row_h = int(h / len(rows))
    cy = y
    for r, row in enumerate(rows):
        cx = x
        for c, val in enumerate(row):
            is_header = r == 0
            if is_header:
                fill = BLUE
                color = WHITE
            elif highlight_col is not None and c == highlight_col:
                fill = SHADE
                color = DARK
            else:
                fill = RGBColor(247, 251, 255) if r % 2 == 0 else WHITE
                color = TEXT
            rect(slide, cx, cy, col_ws[c], row_h, fill=fill, line=GRID)
            text(slide, val, cx + emu(0.04), cy + emu(0.04),
                 col_ws[c] - emu(0.08), row_h - emu(0.06),
                 size=size, color=color, bold=is_header or (highlight_col == c and not is_header),
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            cx += col_ws[c]
        cy += row_h


def extract_figures():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    doc = Document(PAPER)
    body_children = list(doc.element.body.iterchildren())
    figures = {}
    for idx, child in enumerate(body_children):
        if child.tag.split("}")[-1] != "p":
            continue
        blips = list(child.iter(qn("a:blip")))
        if not blips:
            continue
        caption = ""
        for next_child in body_children[idx + 1: idx + 4]:
            if next_child.tag.split("}")[-1] != "p":
                continue
            candidate = "".join(t.text or "" for t in next_child.iter(qn("w:t"))).strip()
            if candidate.startswith("图"):
                caption = candidate
                break
        m = re.search(r"图\s*(\d+)-(\d+)", caption)
        if not m:
            continue
        key = f"{m.group(1)}-{m.group(2)}"
        rid = blips[0].get(qn("r:embed")) or blips[0].get(qn("r:link"))
        if not rid:
            continue
        part = doc.part.related_parts[rid]
        ext = ".png" if part.content_type == "image/png" else ".jpg"
        path = ASSET_DIR / f"fig_{key.replace('-', '_')}{ext}"
        path.write_bytes(part.blob)
        figures[key] = path
    return figures


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------


def cover(slide):
    set_bg(slide, RGBColor(236, 245, 253))
    rect(slide, emu(0), emu(0), emu(3.7), emu(7.5), fill=BLUE, line=BLUE)
    rect(slide, emu(3.68), emu(0), emu(0.18), emu(7.5), fill=ORANGE, line=ORANGE)
    text(slide, "穷究于理  成就于工", emu(0.52), emu(0.55), emu(2.7), emu(0.28),
         size=15, color=WHITE, bold=True)
    text(slide, "本科毕业论文答辩", emu(0.52), emu(1.18), emu(2.4), emu(0.24),
         size=13.5, color=RGBColor(216, 235, 250))
    text(slide, "支持多写的 OLTP 数据库一体机\n页面亲和性关键技术研究",
         emu(4.65), emu(1.48), emu(7.6), emu(1.15), size=30, color=DARK, bold=True)
    text(slide, "用在线亲和把共同访问关系转化为本地访问",
         emu(4.70), emu(3.05), emu(7.2), emu(0.30), size=15, color=GREEN, bold=True)
    bullets(slide,
            ["核心问题:多写带来并行度,也带来远程页面所有权转移",
             "核心想法:让经常共同访问的元组更可能被同一节点本地访问",
             "核心机制:在线采样 → 图划分 → 快照发布 → 后台迁移",
             "核心证据:主对照 + 长稳 + 迁移转化 + 消融闭环"],
            emu(4.85), emu(3.75), emu(4.30), emu(1.20), size=12)
    text(slide, "汇报人:明泰\n专业班级:智能科学与技术4班\n指导教师:温泉、卢卫\n汇报时间:2026年5月",
         emu(9.30), emu(4.12), emu(2.95), emu(0.92), size=11, color=TEXT)
    text(slide, "CDUT", emu(0.55), emu(6.55), emu(1.3), emu(0.32),
         size=18, color=WHITE, bold=True)


def slide_storyline(slide):
    """2 - Defense storyline and headline evidence"""
    set_bg(slide)
    header(slide, "答辩主线:五个问题串起证据链", "00 总览")
    lead(slide,
         "导读:本次答辩按五个问题推进——为何做 / 旧法为何不够 / 机制如何展开 / 正确性如何守住 / 实验如何支撑。每问对应一段证据,贯通形成一条完整的「证据链」。")

    # Left: one-sentence problem and method chain.
    rect(slide, emu(0.55), emu(1.00), emu(5.55), emu(5.30), fill=GRAY, line=GRID)
    text(slide, "中心论点", emu(0.78), emu(1.18), emu(5.1), emu(0.28),
         size=15, color=DARK, bold=True)
    text(slide,
         "多写云数据库中,节点间的页面所有权转移会成为并行处理的隐性瓶颈;本文以事务共同访问关系为输入,在线构建亲和图、计算目标节点、后台迁移元组,把跨节点协调成本系统性降低。",
         emu(0.78), emu(1.58), emu(5.05), emu(1.00),
         size=11.5, color=DARK, bold=True)
    text(slide, "证据链结构", emu(0.78), emu(2.75), emu(5.1), emu(0.24),
         size=13.5, color=GREEN, bold=True)
    steps = [
        ("为何做", "节点间数据传输抵消多写并行收益"),
        ("旧法不足", "静态 / 路由 / Lazy Release 都缺在线学习"),
        ("怎么做", "采样 → 建图 → 划分 → 发布 → 迁移"),
        ("怎么守住", "快照、epoch、线性化点 + 5 条不变量"),
        ("怎么证明", "吞吐 +12% / 远程 −5pp / 6 h 长稳 / 消融"),
    ]
    for i, (tag, body) in enumerate(steps):
        y = 3.12 + i * 0.52
        rect(slide, emu(0.83), emu(y), emu(0.95), emu(0.38), fill=BLUE if i < 3 else GREEN, line=None)
        text(slide, tag, emu(0.83), emu(y + 0.08), emu(0.95), emu(0.18),
             size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, body, emu(1.95), emu(y + 0.06), emu(3.85), emu(0.22),
             size=10.5, color=TEXT)

    # Right: four claims that the rest of the deck proves.
    rect(slide, emu(6.35), emu(1.00), emu(6.43), emu(5.30), fill=WHITE, line=GRID)
    text(slide, "每个问题对应一组证据页", emu(6.58), emu(1.18), emu(5.9), emu(0.28),
         size=15, color=DARK, bold=True)
    claims = [
        ("Q1", "问题是否真实", "图2-1/2-2 说明远程锁、页面推送和所有权转移为何进入路径。"),
        ("Q2", "为什么需要新机制", "三类基线都不能同时满足在线学习、低侵入发布和后台落实。"),
        ("Q3", "机制如何闭环", "50 ms 聚合、10 s 重分区、200 ms 限速迁移组成异步闭环。"),
        ("Q4", "正确性如何守住", "BLink 切换是线性化点,双页 X 锁和 epoch 保护半迁移/过期计划。"),
        ("Q5", "实验如何支撑", "吞吐 +5.16%/+11.79%,远程访问降至 19.660%,6 h 稳定运行。"),
    ]
    for i, (tag, title, body) in enumerate(claims):
        y = 1.58 + i * 0.75
        color = [ORANGE, BLUE, GREEN, DARK, RED][i]
        rect(slide, emu(6.58), emu(y), emu(0.68), emu(0.55), fill=color, line=color)
        text(slide, tag, emu(6.58), emu(y + 0.13), emu(0.68), emu(0.22),
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title, emu(7.42), emu(y - 0.01), emu(1.58), emu(0.22),
             size=11.6, color=color, bold=True)
        text(slide, body, emu(8.92), emu(y - 0.01), emu(3.45), emu(0.46),
             size=9.3, color=TEXT)

    metric(slide, "28912.06", "10 min 在线亲和 TPS", emu(6.70), emu(5.42), emu(1.65), emu(0.60), GREEN)
    metric(slide, "19.660%", "远程访问比例", emu(8.60), emu(5.42), emu(1.55), emu(0.60), BLUE)
    metric(slide, "7.23197 亿", "6 h 总事务数", emu(10.40), emu(5.42), emu(1.72), emu(0.60), GREEN)

    conclusion(slide,
               "讲法按五问推进:先证明问题真实,再证明方案必要、机制闭环、正确性守住、实验支撑。",
               bridge="先回答第一个问题:这个系统里为什么会出现页面亲和性瓶颈。")


def slide_research_object(slide, figs):
    """3 - 研究对象 + 一句话贡献"""
    set_bg(slide)
    header(slide, "研究对象:支持多写的 OLTP 数据库一体机", "01 背景")
    lead(slide,
         "导读:支持多写 = 多台计算服务器 *同时* 处理写事务;一体机 = 软硬件深度集成的数据库系统(对标 PolarDB、GaussDB)。多写带来并发能力,也带来跨节点协调成本——本文优化的就是后者。")
    image_fit(slide, figs.get("2-2"), emu(0.55), emu(1.00), emu(7.55), emu(5.30),
              "图2-2 WookongDB-MP 三层架构")
    rect(slide, emu(8.35), emu(1.00), emu(4.43), emu(2.40), fill=GRAY, line=GRID)
    text(slide, "三层架构", emu(8.55), emu(1.18), emu(4.0), emu(0.26),
         size=14.5, color=DARK, bold=True)
    bullets(slide,
            ["存储层:DiskManager / LogManager / WAL,负责持久化页与日志",
             "远程服务层:PageTable / Partition / TimeStamp / AssignmentTable",
             "计算层:事务执行 + 本地页缓存 + 2PL + 亲和迁移线程"],
            emu(8.55), emu(1.55), emu(4.0), emu(1.75), size=10.5)
    rect(slide, emu(8.35), emu(3.55), emu(4.43), emu(2.75), fill=WHITE, line=GRID)
    text(slide, "本文贡献", emu(8.55), emu(3.73), emu(4.0), emu(0.26),
         size=14.5, color=GREEN, bold=True)
    text(slide, "在 lazy 模式上方增加在线亲和层,把共同访问元组逐步聚到更合适的计算节点。",
         emu(8.55), emu(4.05), emu(4.05), emu(0.95), size=11, color=DARK, bold=True)
    bullets(slide,
            ["不改 Lazy Release / 不改共享存储读写协议",
             "把 shared-nothing 的图划分思想改造成共享存储在线优化",
             "已在 WookongDB-MP 上端到端实现并多机验证"],
            emu(8.55), emu(5.05), emu(4.0), emu(1.20), size=10.5)
    conclusion(slide,
               "研究边界很明确:不重写共享存储协议,只在计算层/远程服务层增加在线亲和控制面。",
               bridge="系统边界清楚后,下一步说明为什么远程路径会抵消多写收益。")


def slide_bottleneck(slide, figs):
    """4 - 瓶颈:Lazy Release 能保留但不能创造"""
    set_bg(slide)
    header(slide, "瓶颈:远程所有权转移如何进入事务关键路径", "01 背景")
    lead(slide,
         "导读:同一时刻同一页面只能由一个节点修改;事务访问的元组若在远端节点持有的页面里,必须先通过「请求 → 等待 → 推送 → 加锁」把页面所有权拿回来——这段时间往往远超事务自身的计算。")
    image_fit(slide, figs.get("2-1"), emu(0.55), emu(1.00), emu(7.55), emu(5.30),
              "图2-1 Lazy Release 页面访问流程")
    rect(slide, emu(8.35), emu(1.00), emu(4.43), emu(2.45), fill=WHITE, line=GRID)
    text(slide, "远程路径开销构成", emu(8.55), emu(1.18), emu(4.0), emu(0.26),
         size=14.5, color=DARK, bold=True)
    bullets(slide,
            ["请求远端锁 → 等待持有者释放",
             "页面推送 + 本地安装",
             "WAL 开启时还要先满足日志持久化约束",
             "网络、锁等待、日志路径都会放大尾延迟"],
            emu(8.55), emu(1.55), emu(4.0), emu(1.80), size=10.5)
    rect(slide, emu(8.35), emu(3.60), emu(4.43), emu(2.70), fill=GRAY, line=GRID)
    text(slide, "关键洞察", emu(8.55), emu(3.78), emu(4.0), emu(0.26),
         size=14.5, color=ORANGE, bold=True)
    text(slide, "Lazy Release 假设 \"近期访问的数据继续被同一节点访问\"——这个假设在亲和负载下并不成立。",
         emu(8.55), emu(4.10), emu(4.05), emu(0.95), size=11, color=DARK, bold=True)
    bullets(slide,
            ["例:转账 A→B,A 在节点 1、B 在节点 2 → 节点 1 须先取回 B 所在页的所有权",
             "Lazy Release 只能保留已有局部性,无法主动汇聚强相关元组",
             "本文优化对象:减少未来进入远程路径的次数,而非替代协议"],
            emu(8.55), emu(5.10), emu(4.0), emu(1.15), size=10.0)
    conclusion(slide,
               "因果链是:相关元组分散 → 事务跨节点取页 → 锁等待/推送/WAL放大尾延迟 → 需要减少远程路径次数。",
               bridge="如果目标是减少远程路径次数,现有三类办法为什么不够?")


def slide_three_approaches(slide):
    """5 - 三种现有方案为什么都不够"""
    set_bg(slide)
    header(slide, "三类现有方案都不足以解决问题", "01 背景")
    lead(slide,
         "导读:把三类已有思路按「是否搬数据 / 是否在线学习」摆在同一张图上看,可以发现它们的共同缺口——没有一个机制能让数据布局 *持续* 跟随运行时访问关系。")
    cols = [
        ("Lazy Release",
         "[ 保留局部性 ]",
         BLUE,
         ["节点持有页面时不主动释放",
          "复用本地缓存 → 后续事务零远程开销",
          "对路由稳定 + 热点稳定的场景效果好"],
         ["不能主动搬动数据,只能保留已有局部性",
          "强相关元组分散在不同节点时无能为力",
          "等待局部性自然出现,亲和负载下不成立"],
         "24.114%", "无亲和基线远程访问比例", "(主对照基线)"),
        ("Schism 静态分区",
         "[ 一次性快照 ]",
         ORANGE,
          ["离线训练事务图 + METIS 划分",
          "2009.3 万事务训练 + 120 s apply 后测量",
          "shared-nothing 数据库的经典思路"],
         ["共享存储 ≠ shared-nothing 假设",
          "训练完成后无法跟随运行时热点漂移",
          "测量阶段 pageIdChanges = 0,没有持续修正能力"],
         "24.875%", "Schism 静态远程访问 ↑", "反而比无亲和升 0.76 pp"),
        ("哈希 / 静态路由",
         "[ 只动事务,不动数据 ]",
         RED,
         ["MP-Router 仅按主键散列",
          "事务入口分散到 4 个计算节点",
          "工程实现最简单,无需训练 / 后台线程"],
         ["不改物理布局,数据仍在原页",
          "强相关元组仍跨节点协调",
          "并发上去后远程比例不下降"],
         "Mode 23", "等价消融:无 ownership 信号", "比 Mode 13 吞吐 −26.4%"),
    ]
    box_w = 4.05
    for i, (name, tag, color, did, missing, val, val_label, val_sub) in enumerate(cols):
        x = emu(0.55 + i * (box_w + 0.18))
        rect(slide, x, emu(1.05), emu(box_w), emu(5.20), fill=WHITE, line=GRID)
        rect(slide, x, emu(1.05), emu(box_w), emu(0.55), fill=color, line=color)
        text(slide, name, x + emu(0.20), emu(1.18), emu(box_w - 0.4), emu(0.30),
             size=15.5, color=WHITE, bold=True)
        text(slide, tag, x + emu(0.20), emu(1.78), emu(box_w - 0.4), emu(0.22),
             size=10, color=color, italic=True, bold=True)
        # 工作方式 section
        text(slide, "工作方式", x + emu(0.20), emu(2.18), emu(box_w - 0.4), emu(0.22),
             size=10.5, color=BLUE, bold=True)
        bullets(slide, did, x + emu(0.22), emu(2.48), emu(box_w - 0.4), emu(1.30),
                size=10.5, spacing=4)
        # 主要局限 section
        text(slide, "主要局限", x + emu(0.20), emu(3.85), emu(box_w - 0.4), emu(0.22),
             size=10.5, color=ORANGE, bold=True)
        bullets(slide, missing, x + emu(0.22), emu(4.15), emu(box_w - 0.4), emu(1.30),
                size=10.5, spacing=4)
        # 实测信号 callout box
        rect(slide, x + emu(0.18), emu(5.45), emu(box_w - 0.36), emu(0.72), fill=SHADE, line=color)
        text(slide, val, x + emu(0.30), emu(5.50), emu(1.30), emu(0.38),
             size=18, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
        text(slide, val_label, x + emu(1.55), emu(5.50), emu(box_w - 1.75), emu(0.22),
             size=9.5, color=DARK, bold=True)
        text(slide, val_sub, x + emu(1.55), emu(5.78), emu(box_w - 1.75), emu(0.22),
             size=8.5, color=MUTED, italic=True)
    conclusion(slide,
               "由此得到设计约束:必须在线学习亲和关系、低侵入发布目标节点、并在后台把目标落实为物理局部性。",
               bridge="进入方案前先澄清三个容易混淆的概念,避免把逻辑分区误认为物理迁移。")


def slide_concept_decoupling(slide):
    """6 - 关键概念解耦"""
    set_bg(slide)
    header(slide, "概念解耦:逻辑/物理位置 + 关键性能代理指标", "02 想法")
    lead(slide,
         "导读:本页澄清两组易混淆的概念。第一组「逻辑/物理/所有者」类比户籍、当前居住地、具体门牌号——三者可独立变化。第二组「EdgeCut/远程访问比例」分别是图论指标和实际性能代理。")

    # Left: 3-layer stack
    rect(slide, emu(0.55), emu(1.05), emu(6.10), emu(5.25), fill=GRAY, line=GRID)
    text(slide, "三层解耦", emu(0.78), emu(1.22), emu(5.6), emu(0.28),
         size=15, color=DARK, bold=True)

    layers = [
        ("逻辑目标节点  assignment(v)", "ParMETIS adaptive 输出  ·  存放在 AssignmentTable",
         GREEN, 1.65),
        ("当前页面所有者", "PageTable + Lazy Release  ·  随事务执行动态变化",
         BLUE, 3.20),
        ("物理 Rid 位置", "BLink 索引 + 共享存储页面 slot  ·  由迁移协议改变",
         ORANGE, 4.75),
    ]
    for label, sub, color, y in layers:
        rect(slide, emu(0.85), emu(y), emu(5.50), emu(1.30), fill=WHITE, line=GRID)
        rect(slide, emu(0.85), emu(y), emu(0.18), emu(1.30), fill=color, line=color)
        text(slide, label, emu(1.18), emu(y + 0.18), emu(5.1), emu(0.30),
             size=13.5, color=DARK, bold=True)
        text(slide, sub, emu(1.18), emu(y + 0.62), emu(5.1), emu(0.50),
             size=10.5, color=MUTED)
    text(slide, "↓ MigrationWorker 把逻辑分配落到物理位置",
         emu(0.95), emu(2.96), emu(5.40), emu(0.20), size=9.5, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)
    text(slide, "↓ Lazy Release 把物理位置反映为当前所有者",
         emu(0.95), emu(4.51), emu(5.40), emu(0.20), size=9.5, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)

    # Right: edgecut vs remote-access metric comparison
    rect(slide, emu(6.85), emu(1.05), emu(5.93), emu(5.25), fill=WHITE, line=GRID)
    text(slide, "为什么用「远程访问比例」而不是 edgecut", emu(7.05), emu(1.22), emu(5.5), emu(0.28),
         size=14.5, color=DARK, bold=True)

    rect(slide, emu(7.05), emu(1.65), emu(5.55), emu(1.45), fill=GRAY, line=GRID)
    text(slide, "edgecut    [ 图论指标 ]", emu(7.25), emu(1.78), emu(5.2), emu(0.28),
         size=13, color=MUTED, bold=True)
    bullets(slide,
            ["ParMETIS 直接输出,衡量切开的亲和边权",
             "不区分远程访问 vs 本地访问的代价差异"],
            emu(7.25), emu(2.15), emu(5.2), emu(0.85), size=10.5, color=MUTED)

    rect(slide, emu(7.05), emu(3.18), emu(5.55), emu(2.95), fill=SHADE, line=GREEN)
    text(slide, "远程访问比例    [ 关键性能代理指标 ]", emu(7.25), emu(3.32), emu(5.2), emu(0.28),
         size=13, color=GREEN, bold=True)
    bullets(slide,
            ["直接对应 Lazy Release 协议的远程锁等待、页面推送、所有权转移",
             "WAL 开启时还包含日志刷盘等待",
             "主实验 / 消融的 headline 指标"],
            emu(7.25), emu(3.70), emu(5.2), emu(1.55), size=11)
    text(slide, "消融周期实验印证:20s 周期 edgecut 由 12 426 升至 44 092,但吞吐略涨 → 两者并不等价。",
         emu(7.25), emu(5.40), emu(5.2), emu(0.60), size=10, color=DARK, italic=True)

    conclusion(slide,
               "这一页回答方案的语义边界:AssignmentTable 是控制面目标,远程访问比例才是最终性能代理。",
               bridge="有了语义边界,机制第一步就是从已提交事务里学习近期亲和关系。")


def slide_observe(slide, figs):
    """7 - 观察:事务采样 + 亲和图建模"""
    set_bg(slide)
    header(slide, "观察:事务采样 + 元组级亲和图建模", "03 机制 1 / 6")
    lead(slide,
         "导读:亲和关系是运行时浮现的——离线训练再充分,也跟不上真实负载的热点漂移。所以第一步必须在线、轻量、连续地观察事务访问关系,把它结构化为一张可被算法消费的图。")
    image_fit(slide, figs.get("3-1"), emu(0.55), emu(1.00), emu(7.45), emu(5.30),
              "图3-1 亲和性元组迁移总体架构")
    rect(slide, emu(8.25), emu(1.00), emu(4.55), emu(5.30), fill=GRAY, line=GRID)
    text(slide, "三段闭环:观察 → 决策 → 执行", emu(8.45), emu(1.18), emu(4.2), emu(0.26),
         size=13.5, color=DARK, bold=True)
    bullets(slide,
            ["SampleRing:只在事务提交后记录 A(T)",
             "Aggregator:50 ms 周期合并为亲和图",
             "中止事务不入图,降低并发冲突带来的噪声"],
            emu(8.45), emu(1.50), emu(4.2), emu(1.55), size=10.5)
    text(slide, "亲和图的两类权重", emu(8.45), emu(3.10), emu(4.2), emu(0.26),
         size=13.5, color=BLUE, bold=True)
    text(slide, "A(T) = {t₁,…,tₖ}   w(tᵢ,tⱼ) += 1   c(tᵢ) += 1",
         emu(8.45), emu(3.45), emu(4.2), emu(0.28), size=12, color=DARK, bold=True,
         align=PP_ALIGN.CENTER)
    bullets(slide,
            ["边权 w:共同访问强度 → 同分区目标",
             "顶点权 c:元组热度 → ParMETIS 负载均衡",
             "边权阈值 + 衰减 λ + TTL 清冷元组,适应热点漂移"],
            emu(8.45), emu(3.80), emu(4.2), emu(2.40), size=10.5)
    conclusion(slide,
               "采样阶段以已提交事务为唯一输入源;单事务边数阈值、衰减因子 λ 与 TTL 共同约束图规模与噪声。",
               bridge="亲和关系已建模,下一步把图论目标翻译为可被 ParMETIS 求解的优化问题。")


def slide_partition_objective(slide, figs):
    """8 - Partition objective and database semantics"""
    set_bg(slide)
    header(slide, "从亲和图到 ParMETIS 划分目标", "03 机制 2 / 6")
    lead(slide,
         "导读:EdgeCut = 把元组分到不同节点时被 \"切断\" 的共同访问边权重之和;越小,高亲和元组就越能保留在同一节点。本页给出这一目标如何映射到数据库语义,以及在工程上做出的三个折中。")
    image_fit(slide, figs.get("3-2"), emu(0.55), emu(1.00), emu(6.65), emu(5.30),
              "图3-2 ParMETIS 亲和图划分示意")

    rect(slide, emu(7.45), emu(1.00), emu(5.33), emu(2.38), fill=GRAY, line=GRID)
    text(slide, "图划分目标如何落到数据库", emu(7.68), emu(1.18), emu(4.9), emu(0.28),
         size=14.5, color=DARK, bold=True)
    bullets(slide,
            ["V:元组;E:同一提交事务共同访问形成的带权边",
             "EdgeCut 越低 → 高亲和元组越少被分到不同节点",
             "顶点权和 ubvec 控制负载,避免把热点集中到少数节点"],
            emu(7.68), emu(1.58), emu(4.90), emu(1.55), size=10.7)

    rect(slide, emu(7.45), emu(3.58), emu(5.33), emu(2.72), fill=WHITE, line=GRID)
    text(slide, "三个工程取舍", emu(7.68), emu(3.76), emu(4.9), emu(0.26),
         size=14.5, color=BLUE, bold=True)
    rows = [
        ("普通图而非超图",
         "事务访问 k 个元组时展开为 k(k−1)/2 条边;SmallBank 单事务访问元组少,近似代价可接受。"),
        ("adaptive repartition",
         "上一轮 assignment 作为参考,结合 maxChangedVerticesRatio 控制过大迁移。"),
        ("指标边界",
         "EdgeCut 是逻辑图指标;最终仍以远程访问比例衡量实际页面路径收益。"),
    ]
    for i, (title, body) in enumerate(rows):
        y = 4.16 + i * 0.62
        text(slide, title, emu(7.68), emu(y), emu(1.75), emu(0.22),
             size=10.5, color=GREEN if i == 1 else ORANGE if i == 2 else BLUE, bold=True)
        text(slide, body, emu(9.38), emu(y - 0.02), emu(3.10), emu(0.42),
             size=9.6, color=TEXT)

    conclusion(slide,
               "划分目标不是单纯最低 EdgeCut,而是在低割、负载均衡和可承受迁移量之间折中。",
               bridge="目标函数讲清楚后,再说明多节点环境里如何一致地产生和分发结果。")


def slide_decide(slide, figs):
    """9 - 决策:分布式图划分 + ParMETIS adaptive", "03 机制 3 / 6"""
    set_bg(slide)
    header(slide, "决策:分布式图划分 + ParMETIS adaptive", "03 机制 3 / 6")
    lead(slide,
         "导读:ParMETIS 是成熟的并行图划分库;\"adaptive\" 意为以上一轮分区为起点做增量调整,避免每次从零划分导致大规模迁移抖动。多节点协同分三步:边交换 → inventory → 分发。")
    image_fit(slide, figs.get("3-3"), emu(0.55), emu(1.00), emu(7.45), emu(5.30),
              "图3-3 分布式图划分流程")
    rect(slide, emu(8.25), emu(1.00), emu(4.55), emu(5.30), fill=GRAY, line=GRID)
    text(slide, "三阶段流程", emu(8.45), emu(1.18), emu(4.2), emu(0.26),
         size=13.5, color=DARK, bold=True)
    bullets(slide,
            ["边交换:按 tupleId 稳定哈希,负责节点合并同边权",
             "Inventory 交换:统一 vtxdist 顶点编号",
             "分区分发:各节点只广播自己负责的 assignment slice"],
            emu(8.45), emu(1.50), emu(4.2), emu(1.65), size=10.5)
    text(slide, "工程关键", emu(8.45), emu(3.20), emu(4.2), emu(0.26),
         size=13.5, color=BLUE, bold=True)
    bullets(slide,
            ["epoch barrier 保证各节点同步同一版本",
             "上一轮分区作为 adaptive 输入 → 降低迁移抖动",
             "ParMETIS = 旁路 MPI 进程,UDS 与主进程通信",
             "主进程不管理 MPI 生命周期,分区失败只影响本轮优化"],
            emu(8.45), emu(3.55), emu(4.2), emu(2.65), size=10.5)
    conclusion(slide,
               "多节点环境下,通过 epoch barrier 同步、按 tupleId 哈希交换边、slice 分发,产出全局一致的分区结果。",
               bridge="分区结果不能直接进入事务路径阻塞前台,因此需要快照式发布。")


def slide_publish(slide):
    """10 - 发布:AssignmentTable 快照 + epoch"""
    set_bg(slide)
    header(slide, "发布:快照式 AssignmentTable + epoch 版本号", "03 机制 4 / 6")
    lead(slide,
         "导读:快照 = 一份不可变的版本;读者读到的就是某一具体版本,期间不会被改写。epoch = 版本号,每次重分区 +1。读路径用原子指针读快照,写路径在后台构建好新快照后整体切换。")

    # ---- Left panel: algorithm + lifecycle ----
    rect(slide, emu(0.55), emu(1.05), emu(6.30), emu(5.25), fill=GRAY, line=GRID)
    text(slide, "算法 3-2  AssignmentTable 查找", emu(0.78), emu(1.18), emu(5.8), emu(0.26),
         size=13.5, color=DARK, bold=True)
    text(slide, "Input:  tupleId, fallbackNode      Output: 元组对应的目标计算节点 nodeId",
         emu(0.78), emu(1.50), emu(5.85), emu(0.24), size=9.5, color=MUTED)

    # Code box (compact)
    code_lines = [
        "1:  snapshot ← currentAssignment              # 原子指针读",
        "2:  if tupleId ∈ snapshot:",
        "3:      return snapshot[tupleId]",
        "4:  return fallbackNode                       # 冷元组退路",
    ]
    rect(slide, emu(0.78), emu(1.84), emu(5.85), emu(1.85), fill=WHITE, line=GRID)
    for i, line in enumerate(code_lines):
        text(slide, line, emu(0.95), emu(1.96 + i * 0.40), emu(5.5), emu(0.30),
             size=10.5, color=DARK)

    # Snapshot lifecycle flow (NEW) — 4 chevrons
    text(slide, "快照生命周期(epoch 单调递增)",
         emu(0.78), emu(3.82), emu(5.85), emu(0.22), size=10.5, color=DARK, bold=True)
    flow = [
        ("ParMETIS\n输出新分区", "epoch += 1", GREEN),
        ("Aggregator\n合并新快照", "build candidate", BLUE),
        ("原子指针\n切换", "compare-and-swap", ORANGE),
        ("读 / 迁移线程\n消费", "snapshot 只读", DARK),
    ]
    step_w = 1.30
    step_h = 1.00
    gap_w = 0.18
    for i, (label, sub, color) in enumerate(flow):
        x = 0.78 + i * (step_w + gap_w)
        rect(slide, emu(x), emu(4.10), emu(step_w), emu(step_h), fill=WHITE, line=color)
        text(slide, label, emu(x + 0.05), emu(4.18), emu(step_w - 0.10), emu(0.50),
             size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, sub, emu(x + 0.05), emu(4.68), emu(step_w - 0.10), emu(0.30),
             size=8.2, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            ax = x + step_w + 0.01
            text(slide, "▶", emu(ax), emu(4.48), emu(gap_w), emu(0.24),
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # 3 italic takeaways below lifecycle
    rect(slide, emu(0.78), emu(5.28), emu(5.85), emu(0.95), fill=WHITE, line=GRID)
    takeaways = [
        "原子指针替换:构建完成后切换,读者只持 snapshot,不需要全局锁",
        "每条记录携带 epoch:迁移线程执行前重新校验是否过期",
        "fallback = 系统原页面归属 / 哈希归属 → 冷元组按既有路径访问",
    ]
    for i, line in enumerate(takeaways):
        text(slide, "  • " + line, emu(0.85), emu(5.34 + i * 0.30), emu(5.75), emu(0.26),
             size=9.5, color=GREEN, italic=True, bold=True)

    # ---- Right: 3 property cards (extended) ----
    cards = [
        ("低侵入的发布", BLUE,
         "读路径无锁,前台事务不等待分区器;后台合并完成后原子切换 currentAssignment 指针。",
         "fallback 路径处理未进入图的冷元组,不影响访问正确性。"),
        ("逻辑控制面", GREEN,
         "AssignmentTable 表示逻辑目标节点,不直接表示物理位置或当前页面所有者。",
         "三层(逻辑 / 所有者 / 物理)解耦,使 ParMETIS 频繁更新不绑死实际数据位置。"),
        ("过期保护", ORANGE,
         "ParMETIS 每轮分配 epoch ↑;迁移线程执行前重新解析 Rid 并校验版本。",
         "失败任务不进入冷却,后续新 assignment 可立即重新规划该元组。"),
    ]
    for i, (title, color, body1, body2) in enumerate(cards):
        y = 1.05 + i * 1.78
        rect(slide, emu(7.05), emu(y), emu(5.75), emu(1.62), fill=WHITE, line=GRID)
        rect(slide, emu(7.05), emu(y), emu(0.18), emu(1.62), fill=color, line=color)
        text(slide, title, emu(7.35), emu(y + 0.15), emu(5.4), emu(0.28),
             size=14, color=color, bold=True)
        text(slide, body1, emu(7.35), emu(y + 0.50), emu(5.4), emu(0.55),
             size=11, color=TEXT)
        text(slide, "▸ " + body2, emu(7.35), emu(y + 1.10), emu(5.4), emu(0.45),
             size=10, color=MUTED, italic=True)

    conclusion(slide,
               "AssignmentTable 以原子指针切换实现低侵入发布:读路径无锁,写路径后台构建后整体生效,前台事务零阻塞。",
               bridge="发布仅传递「应该在哪」;真正形成物理局部性须由后台迁移安全落实。")


def slide_execute(slide, figs):
    """11 - 执行:元组迁移协议"""
    set_bg(slide)
    header(slide, "执行:元组迁移协议(BLink 切换 = 线性化点)", "03 机制 5 / 6")
    lead(slide,
         "导读:线性化点 = 整个迁移对外可见的那个瞬间。BLink 索引更新前,前台事务通过索引只能到达源位置;更新后,只能到达目标位置——不存在「同一元组两份可见副本」的中间状态。")
    image_fit(slide, figs.get("3-4"), emu(0.55), emu(1.00), emu(7.45), emu(5.30),
              "图3-4 元组迁移协议流程")
    rect(slide, emu(8.25), emu(1.00), emu(4.55), emu(5.30), fill=GRAY, line=GRID)
    text(slide, "7 步执行 + 线性化点", emu(8.45), emu(1.18), emu(4.2), emu(0.26),
         size=13.5, color=DARK, bold=True)
    bullets(slide,
            ["BLink 定位元组当前 srcPage / slot",
             "重新解析 Rid,校验源页是否仍属本节点",
             "按全局页号序取双 X 锁 → 避免循环等待",
             "复制 key + DataItem + value 到目标 slot",
             "更新 BLink → 线性化点",
             "清源 slot + 写 insert/delete 日志",
             "释放锁并更新统计"],
            emu(8.45), emu(1.50), emu(4.2), emu(2.85), size=10.5)
    text(slide, "失败与跳过", emu(8.45), emu(4.45), emu(4.2), emu(0.26),
         size=13.5, color=ORANGE, bold=True)
    bullets(slide,
            ["目标 slot 已写但 BLink 未切 → 回滚目标 slot",
             "DataItem 仍被前台事务持锁 → 跳过,不阻塞业务事务"],
            emu(8.45), emu(4.80), emu(4.2), emu(1.40), size=10.5)
    conclusion(slide,
               "BLink 索引更新作为迁移线性化点:前台事务在切换前后均观察到单一可见位置,不暴露中间状态。",
               bridge="协议流程已给出,下面将正确性约束压缩为可审查的五条不变量。")


def slide_correctness(slide):
    """12 - 正确性:5 条不变量"""
    set_bg(slide)
    header(slide, "正确性:5 条不变量", "03 机制 6 / 6")
    lead(slide,
         "导读:不变量 = 不论何时介入、不论以何种顺序观察,都必须始终成立的性质。本页把迁移正确性压缩为五条可逐一审查的不变量,覆盖可见性、并发控制、失败回滚与版本一致性。")
    rows = [
        ["编号", "不变量", "保证手段"],
        ["I1", "同一元组不能同时存在两个可见位置",
         "BLink 索引更新作为迁移线性化点"],
        ["I2", "前台事务不能观察到半迁移状态",
         "源页 + 目标页 双 X 锁保护复制 / 切换 / 清理"],
        ["I3", "迁移失败可回滚或重试",
         "BLink 切换前失败 → 回滚目标 slot;in-flight 状态释放"],
        ["I4", "不破坏元组锁语义",
         "迁移前检查 DataItem 是否 UNLOCKED,否则跳过"],
        ["I5", "过期计划不能错误迁移",
         "执行前重新解析 Rid + 校验 assignment epoch"],
    ]
    table(slide, rows, emu(0.55), emu(1.05), emu(12.23), emu(4.55),
          widths=[0.65, 3.0, 4.7], size=10.5)
    rect(slide, emu(0.55), emu(5.75), emu(12.23), emu(0.55), fill=SHADE, line=GREEN)
    text(slide, "实现要点:复用页面锁接口、BLink 与 WAL 语义,把迁移限制为后台优化而非新事务协议。",
         emu(0.78), emu(5.85), emu(11.77), emu(0.30), size=11.5, color=DARK, bold=True,
         align=PP_ALIGN.CENTER)
    conclusion(slide,
               "五条不变量覆盖可见性、并发控制、失败回滚与版本一致性,将迁移约束为不改变前台事务语义的后台优化。",
               bridge="正确性守住后,仍需说明迁移线程为何不会对前台事务路径造成性能干扰。")


def slide_batch(slide, figs):
    """13 - 性能优化:批量迁移"""
    set_bg(slide)
    header(slide, "性能优化:按 (tableId, srcPage, dstNode) 合并批量迁移", "04 实现")
    lead(slide,
         "导读:页级开销 = 每次访问页面都要做的固定动作(取页锁、定位空闲槽、维护元数据),与移动的元组数量无关。若同一源页的多个元组都要去同一目标节点,可一次性合并迁移,把固定开销摊到多个元组上。")
    image_fit(slide, figs.get("4-1"), emu(0.55), emu(1.00), emu(7.10), emu(5.30),
              "图4-1 批量迁移核心流程")
    rect(slide, emu(7.90), emu(1.00), emu(4.88), emu(2.50), fill=GRAY, line=GRID)
    text(slide, "为什么需要批量", emu(8.10), emu(1.18), emu(4.5), emu(0.26),
         size=13.5, color=DARK, bold=True)
    bullets(slide,
            ["单 tuple 迁移会反复定位源页 / 取源页锁 / 分配目标页",
             "这些是「页级」固定开销,与 tuple 数量无关",
             "同源页 + 同目标的计划可合并,一次锁多 tuple"],
            emu(8.10), emu(1.55), emu(4.5), emu(1.85), size=10.5)
    rect(slide, emu(7.90), emu(3.60), emu(4.88), emu(2.70), fill=WHITE, line=GRID)
    text(slide, "合并粒度的边界", emu(8.10), emu(3.78), emu(4.5), emu(0.26),
         size=13.5, color=BLUE, bold=True)
    bullets(slide,
            ["仅按 (table, srcPage, dstNode) 三元组合并",
             "不跨源页合并 → 控制锁持有时间和死锁风险",
             "组内每元组仍单独校验 BLink / DataItem / epoch",
             "目标页池复用空闲 slot,摊销 updatePageSpace 元数据维护"],
            emu(8.10), emu(4.15), emu(4.5), emu(2.10), size=10.5)
    conclusion(slide,
               "批量迁移以页级开销摊销为主要收益:前台吞吐基本持平(+0.04%),积压下降 35.37%、失败次数下降 77.78%。",
               bridge="机制链条已完整;下面以四个研究问题逐一回扣前述主张。")


def slide_experiment_setup(slide):
    """14 - 实验设置 + 4 RQ + 诚实标注"""
    set_bg(slide)
    header(slide, "实验设置与四个研究问题", "05 实验")
    lead(slide,
         "导读:实验设计的目的不是堆数字,而是让每一个 RQ 独立检验机制的一个环节——能否提升(RQ1)、能否长稳(RQ2)、是否真转化(RQ3)、哪些设计关键(RQ4)。口径条件全部公开标注,不绕开。")

    # Top: topology + key params (left), 4 RQ (right)
    rect(slide, emu(0.55), emu(1.05), emu(6.30), emu(4.55), fill=GRAY, line=GRID)
    text(slide, "拓扑与关键参数", emu(0.78), emu(1.22), emu(5.8), emu(0.26),
         size=14, color=DARK, bold=True)
    bullets(slide,
            ["5 台服务器:4 计算节点 + 1 存储/服务节点",
             "Hygon C86 5390 16C(node0 vCPU/内存 ≈ 其他节点 2×)",
             "Ubuntu 22.04.5 + GCC 11.4 + brpc + ParMETIS",
             "数据来自 summary.txt / MP-Router 日志 / affinityTimeseries.csv"],
            emu(0.78), emu(1.58), emu(5.8), emu(1.60), size=10.5)
    rows = [
        ["参数", "取值", "含义"],
        ["computeHosts × workerThreads", "4 × 32", "总并发约 128"],
        ["systemMode", "13", "ownership + ParMETIS 双信号路由"],
        ["partitionCycleMs", "10 000", "ParMETIS 重分区周期"],
        ["migrationTickMs / batch", "200 / 200", "后台迁移限速"],
        ["decay-factor", "0.9", "亲和图衰减因子"],
        ["tryCount", "157 000", "约 2000 万事务 / 轮"],
    ]
    table(slide, rows, emu(0.78), emu(3.35), emu(5.85), emu(2.15),
          widths=[2.4, 1.0, 2.7], size=8.6)

    rect(slide, emu(7.05), emu(1.05), emu(5.73), emu(4.55), fill=WHITE, line=GRID)
    text(slide, "四个研究问题 = 四个证据关口", emu(7.28), emu(1.22), emu(5.3), emu(0.26),
         size=14, color=BLUE, bold=True)
    rqs = [
        ("RQ1", "在线亲和能否提升吞吐量与尾延迟", "10 min 同条件 vs 无亲和 / Schism"),
        ("RQ2", "在线机制能否长时间稳定运行",      "6 h 长稳实验 7.23 亿事务"),
        ("RQ3", "迁移是否真转化为访问局部性",      "迁移完成率 vs 远程访问比例"),
        ("RQ4", "哪些设计选择是关键",            "批量 / 路由 / 衰减 / 周期 四组消融"),
    ]
    for i, (tag, q, evidence) in enumerate(rqs):
        y = 1.65 + i * 0.92
        rect(slide, emu(7.28), emu(y), emu(0.85), emu(0.78), fill=BLUE, line=BLUE)
        text(slide, tag, emu(7.28), emu(y + 0.18), emu(0.85), emu(0.32),
             size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, q, emu(8.25), emu(y + 0.08), emu(4.4), emu(0.30),
             size=12, color=DARK, bold=True)
        text(slide, evidence, emu(8.25), emu(y + 0.43), emu(4.4), emu(0.28),
             size=10, color=MUTED)

    # Bottom honest caveats
    rect(slide, emu(0.55), emu(5.78), emu(12.23), emu(0.55), fill=SHADE, line=ORANGE)
    text(slide, "口径声明",
         emu(0.78), emu(5.83), emu(1.0), emu(0.42),
         size=10.5, color=ORANGE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    text(slide,
         "WAL 关闭(隔离页面局部性收益)   ·   SmallBank-Aff 亲和负载   ·   Schism 同等训练预算 2009.3 万事务 + 120 s apply   ·   单 run,尚未报告置信区间",
         emu(1.85), emu(5.83), emu(10.85), emu(0.42),
         size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)

    conclusion(slide,
               "实验组织遵循证据关口原则:RQ1-RQ4 分别检验有效性、稳定性、因果转化与设计选择的可分离贡献。",
               bridge="先看 RQ1:同条件 10 min 主对照下,在线亲和的指标改善情况。")


def slide_rq1_charts(slide, figs):
    """15 - RQ1 吞吐 + 尾延迟 图表"""
    set_bg(slide)
    header(slide, "RQ1:在线亲和提升吞吐与尾延迟", "05 实验 RQ1")
    lead(slide,
         "导读:在同一硬件、同一负载、同一时长(10 min)下对比三种方案——无亲和、Schism 静态分区、本文方法。Schism 已给予公平的同等训练预算,排除「基线被故意削弱」的质疑。")
    image_fit(slide, figs.get("5-1"), emu(0.55), emu(1.00), emu(6.10), emu(4.55),
              "图5-1 吞吐结果对比")
    image_fit(slide, figs.get("5-2"), emu(6.78), emu(1.00), emu(6.00), emu(4.55),
              "图5-2 三组方案尾延迟对比")
    metric(slide, "28912.06", "在线亲和 TPS", emu(0.85), emu(5.72), emu(1.85), emu(0.62), GREEN)
    metric(slide, "+11.79%",  "vs Schism", emu(2.85), emu(5.72), emu(1.50), emu(0.62), GREEN)
    metric(slide, "+5.16%",   "vs 无亲和", emu(4.50), emu(5.72), emu(1.50), emu(0.62), GREEN)
    metric(slide, "24.69%",   "Schism P99 高于本方案", emu(6.85), emu(5.72), emu(2.30), emu(0.62), ORANGE)
    metric(slide, "−5.21 pp", "远程访问 vs Schism", emu(9.30), emu(5.72), emu(2.10), emu(0.62), GREEN)
    metric(slide, "19.660%",  "远程访问比例", emu(11.55), emu(5.72), emu(1.30), emu(0.62), BLUE)
    conclusion(slide,
               "RQ1 论证逻辑:若远程路径确为关键瓶颈,远程访问比例的下降应同步反映为吞吐提升与端到端尾延迟收敛——三项指标在主对照下同向改善。",
               bridge="图表给出趋势,下一页用完整数据表呈现口径与 pageId 变更证据。")


def slide_rq1_table(slide):
    """16 - RQ1 主对照数据表"""
    set_bg(slide)
    header(slide, "RQ1:主对照详细数据(同条件 10 min)", "05 实验 RQ1")
    lead(slide,
         "导读:每行一项度量、每列一种方案;高亮列为本方案,与左两列直接相减即得改善幅度。最后两行 pageId 变更与完成迁移数说明本方案在测量阶段持续修正布局,而对照方案在测量阶段完全不动。")
    rows = [
        ["指标", "无亲和", "Schism 静态", "在线亲和(本方案)"],
        ["吞吐量(txn/s)",          "27 492.58", "25 862.32", "28 912.06"],
        ["Warmup 后吞吐(txn/s)",    "27 504.37", "25 872.68", "28 924.78"],
        ["执行 P99(ms)",            "10.40",     "12.03",     "10.28"],
        ["端到端 P99(ms)",          "5 656.14",  "6 735.30",  "5 401.75"],
        ["远程访问比例",            "24.114%",   "24.875%",   "19.660%"],
        ["本地访问比例",            "75.853%",   "75.090%",   "80.303%"],
        ["pageId 变更次数",         "0",         "0",         "117 294"],
        ["完成迁移数",              "0",         "44 200",    "106 049"],
    ]
    table(slide, rows, emu(0.55), emu(1.05), emu(12.23), emu(4.85),
          widths=[2.55, 1.5, 1.7, 2.0], size=10.0, highlight_col=3)
    rect(slide, emu(0.55), emu(6.00), emu(12.23), emu(0.36), fill=SHADE, line=GRID)
    text(slide, "Schism 的端到端 P99 是本方案的 1.2469 倍   ·   在线方案 pageId 变更 117 294 次,说明收益来自运行时持续修正。",
         emu(0.78), emu(6.03), emu(11.77), emu(0.30), size=11, color=DARK, bold=True,
         align=PP_ALIGN.CENTER, italic=False)
    conclusion(slide,
               "主对照闭环:pageId 持续变更 → 远程比例下降 → 吞吐和端到端尾延迟同步改善。",
               bridge="短时有效还不够,下一页回答在线机制是否能长时间稳定运行。")


def slide_rq2(slide, figs):
    """17 - RQ2 6 小时长稳"""
    set_bg(slide)
    header(slide, "RQ2:6 小时长时间稳定性", "05 实验 RQ2")
    lead(slide,
         "导读:把同一套机制连续运行 6 小时,观察吞吐是否随采样、重分区与后台迁移的持续运行出现累积衰减。若机制本身有隐藏缺陷,长时间运行会暴露;吞吐曲线持续平稳,即证机制稳定。")
    image_fit(slide, figs.get("5-3"), emu(0.55), emu(1.00), emu(7.45), emu(5.30),
              "图5-3 6 小时稳定性实验 TPS")
    rect(slide, emu(8.25), emu(1.00), emu(4.55), emu(5.30), fill=GRAY, line=GRID)
    text(slide, "核心统计", emu(8.45), emu(1.18), emu(4.2), emu(0.26),
         size=13.5, color=DARK, bold=True)
    stats = [
        ("7.23197 亿",     "总事务数"),
        ("34 356.98 txn/s", "平均吞吐量"),
        ("21 049.5 s",     "测量时间"),
        ("6.53%",          "TPS 变异系数"),
        ("27 746.27 txn/s","5 min 滚动最小值"),
    ]
    for i, (val, label) in enumerate(stats):
        y = 1.55 + i * 0.88
        rect(slide, emu(8.45), emu(y), emu(4.2), emu(0.78), fill=WHITE, line=GRID)
        text(slide, val, emu(8.55), emu(y + 0.10), emu(2.5), emu(0.36),
             size=15.5, color=GREEN, bold=True)
        text(slide, label, emu(8.55), emu(y + 0.48), emu(4.0), emu(0.24),
             size=10.5, color=MUTED)
    conclusion(slide,
               "6 小时实验中,采样、AssignmentTable 更新与后台迁移线程未引发吞吐持续下滑,验证在线机制具备长时间运行能力。",
               bridge="稳定性仅是必要条件;下一页检验迁移是否真转化为局部性。")


def slide_rq3_counter(slide, figs):
    """18 - RQ3 反直觉 punchline"""
    set_bg(slide)
    header(slide, "RQ3:完成更多迁移 ≠ 更好局部性", "05 实验 RQ3")
    lead(slide,
         "导读:本页是反直觉对照——Schism 把规划好的迁移做得 *更彻底*(完成率更高、失败为零),但最终远程访问比例 *反而更差*。这恰恰说明:局部性目标不能简单等同于迁移数量。")

    # Counterintuitive headline strip
    rect(slide, emu(0.55), emu(1.00), emu(12.23), emu(1.30), fill=SHADE, line=ORANGE)
    text(slide, "反直觉对照", emu(0.78), emu(1.10), emu(2.0), emu(0.22),
         size=11, color=ORANGE, bold=True)

    # Two compare blocks side by side
    rect(slide, emu(0.85), emu(1.36), emu(5.55), emu(0.85), fill=WHITE, line=GRID)
    text(slide, "迁移完成率", emu(1.05), emu(1.42), emu(2.2), emu(0.22),
         size=10, color=MUTED)
    text(slide, "Schism 99.597%", emu(1.05), emu(1.66), emu(2.5), emu(0.30),
         size=14.5, color=BLUE, bold=True)
    text(slide, "↑ 高", emu(3.50), emu(1.66), emu(0.8), emu(0.30),
         size=12.5, color=BLUE, bold=True)
    text(slide, "本方案 99.070%", emu(4.30), emu(1.66), emu(2.1), emu(0.30),
         size=14.5, color=MUTED, bold=True)

    rect(slide, emu(6.95), emu(1.36), emu(5.55), emu(0.85), fill=WHITE, line=GRID)
    text(slide, "远程访问比例 → 局部性结果", emu(7.15), emu(1.42), emu(3.5), emu(0.22),
         size=10, color=MUTED)
    text(slide, "Schism 24.875%", emu(7.15), emu(1.66), emu(2.5), emu(0.30),
         size=14.5, color=ORANGE, bold=True)
    text(slide, "↓ 反过来", emu(9.60), emu(1.66), emu(1.3), emu(0.30),
         size=12.5, color=ORANGE, bold=True)
    text(slide, "本方案 19.660%", emu(10.90), emu(1.66), emu(2.2), emu(0.30),
         size=14.5, color=GREEN, bold=True)

    # Three small images
    image_fit(slide, figs.get("5-4"), emu(0.55), emu(2.45), emu(4.12), emu(3.10),
              "图5-4 迁移任务执行情况")
    image_fit(slide, figs.get("5-5"), emu(4.75), emu(2.45), emu(4.12), emu(3.10),
              "图5-5 本地/远程访问比例随时间变化")
    image_fit(slide, figs.get("5-6"), emu(8.93), emu(2.45), emu(3.85), emu(3.10),
              "图5-6 各计算节点远程访问比例对比")

    # Per-image quick reads
    text(slide, "Schism 计划 44 379 / 完成 44 200,失败 0",
         emu(0.65), emu(5.60), emu(4.00), emu(0.20), size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    text(slide, "本方案 local 持续上升,稳定 ≈ 80%",
         emu(4.85), emu(5.60), emu(4.00), emu(0.20), size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    text(slide, "4 个计算节点同步改善,排除单节点偶然",
         emu(9.00), emu(5.60), emu(3.75), emu(0.20), size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    conclusion(slide,
               "迁移数量并非局部性目标本身;真正的因果传导体现在远程访问比例的下降与本地访问比例的稳态上升。",
               bridge="因果转化已确认,最后剖析哪些设计选择支撑了这条转化通路。")


def slide_rq4_ablation(slide, figs):
    """19 - RQ4 消融四件套(2×2 with fig_5_9 added)"""
    set_bg(slide)
    header(slide, "RQ4:消融实验 — 批量 / 路由 / 衰减 / 周期", "05 实验 RQ4")
    lead(slide,
         "导读:消融 = 把某个模块单独关掉或换设置,看对结果的影响有多大。如果关掉某个模块后机制立刻失效,说明该模块就是核心贡献者;如果几乎无影响,说明它只是细节优化。")

    grid_w = 6.05
    grid_h = 2.45
    quads = [
        ("批量迁移开关",         figs.get("5-7"),  "图5-7",
         "积压 1541 → 996  (−35.37%)   ·   失败 9 → 2",   GREEN, 0.55, 1.00),
        ("路由模式  benefit1=ownership / benefit2=ParMETIS", figs.get("5-8"),  "图5-8",
         "Mode 13(b1+b2)比 Mode 23(仅 b2)吞吐 +35.86%",       BLUE,  6.72, 1.00),
        ("亲和图衰减因子 λ",      figs.get("5-9"),  "图5-9",
         "λ ∈ [0.7, 0.9] 吞吐近持平,尾延迟差异明显;结合长稳选择 0.9",   ORANGE,0.55, 3.65),
        ("重分区周期",            figs.get("5-10"), "图5-10",
         "10 s edgecut 12 426  ·  20 s 升至 44 092,吞吐略涨",   RED,   6.72, 3.65),
    ]
    for title, fig, fig_name, finding, color, x, y in quads:
        rect(slide, emu(x), emu(y), emu(grid_w), emu(grid_h), fill=WHITE, line=GRID)
        rect(slide, emu(x), emu(y), emu(0.18), emu(grid_h), fill=color, line=color)
        text(slide, title, emu(x + 0.30), emu(y + 0.06), emu(grid_w - 0.40), emu(0.26),
             size=11.5, color=color, bold=True)
        # image inside the box (top-right region)
        image_fit(slide, fig, emu(x + 0.30), emu(y + 0.38), emu(grid_w - 0.50), emu(1.55), fig_name)
        text(slide, finding, emu(x + 0.30), emu(y + 2.00), emu(grid_w - 0.50), emu(0.40),
             size=10, color=TEXT, italic=True)

    conclusion(slide,
               "性能收益来自路由、亲和分区、后台迁移与关键参数的协同贡献——任一模块单独优化均不构成完整解释。",
               bridge="最后回到五问主线,总结已证明的部分与仍需承认的边界。")


def slide_summary(slide):
    """20 - 局限 + 总结 + 致谢"""
    set_bg(slide, RGBColor(247, 250, 254))
    rect(slide, emu(0), emu(0), emu(13.333), emu(0.72), fill=BLUE, line=BLUE)
    text(slide, "局限、工作总结与未来展望", emu(0.55), emu(0.22), emu(12.0), emu(0.32),
         size=21, color=WHITE, bold=True)

    cols = [
        ("答辩边界   [ 哪些还没证明 ]", ORANGE,
         ["实验主负载是 SmallBank-Aff; YCSB / TPC-C / 动态热点仍需扩展",
          "主性能实验关闭 WAL,结论聚焦页面局部性收益",
          "迁移积压尚未消除,分区生成 ↔ 后台执行仍需平衡",
          "BLink + 迁移日志恢复语义还需系统故障注入验证"]),
        ("故事回扣   [ 五问已回答 ]", GREEN,
         ["问题真实:远程页面路径进入事务关键路径",
          "方案必要:静态 / Lazy / 纯路由都缺在线闭环",
          "机制闭环:采样建图、分区发布、后台迁移",
          "正确性守住:BLink 线性化点 + epoch + 5 条不变量",
          "实验支撑:主对照、6 h 长稳、局部性转化、消融"]),
        ("未来展望", BLUE,
         ["利用 AssignmentTable 优化事务入口路由",
          "按热度 + 边权 + 锁冲突排序迁移优先级",
          "热点页面更细粒度协调,降低 waitLockSuccess / 页面推送",
          "结合页级重组,把同亲和组元组聚集到更少页面",
          "扩展 YCSB / TPC-C / 动态热点负载验证"]),
    ]
    col_w = 4.10
    for i, (title, color, items) in enumerate(cols):
        x = emu(0.50 + i * (col_w + 0.18))
        rect(slide, x, emu(1.10), emu(col_w), emu(4.95), fill=WHITE, line=GRID)
        rect(slide, x, emu(1.10), emu(col_w), emu(0.45), fill=color, line=color)
        text(slide, title, x + emu(0.22), emu(1.20), emu(col_w - 0.4), emu(0.30),
             size=12.5, color=WHITE, bold=True)
        bullets(slide, items, x + emu(0.22), emu(1.75), emu(col_w - 0.4), emu(4.18),
                size=11.0, spacing=6)

    rect(slide, emu(0), emu(6.30), emu(13.333), emu(1.20), fill=BLUE, line=BLUE)
    text(slide, "请各位老师批评指正",
         emu(0.5), emu(6.55), emu(12.33), emu(0.65),
         size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# -----------------------------------------------------------------------------
# Main
# -----------------------------------------------------------------------------

EXPECTED_COUNT = 20


def main() -> None:
    figs = extract_figures()
    prs = Presentation(TEMPLATE)
    # Trim or grow to EXPECTED_COUNT
    while len(prs.slides) > EXPECTED_COUNT:
        delete_slide(prs, len(prs.slides) - 1)
    while len(prs.slides) < EXPECTED_COUNT:
        add_blank_slide(prs)

    builders = [
        cover,                                            # 1
        slide_storyline,                                  # 2
        lambda s: slide_research_object(s, figs),         # 3
        lambda s: slide_bottleneck(s, figs),              # 4
        slide_three_approaches,                           # 5
        slide_concept_decoupling,                         # 6
        lambda s: slide_observe(s, figs),                 # 7
        lambda s: slide_partition_objective(s, figs),     # 8
        lambda s: slide_decide(s, figs),                  # 9
        slide_publish,                                    # 10
        lambda s: slide_execute(s, figs),                 # 11
        slide_correctness,                                # 12
        lambda s: slide_batch(s, figs),                   # 13
        slide_experiment_setup,                           # 14
        lambda s: slide_rq1_charts(s, figs),              # 15
        slide_rq1_table,                                  # 16
        lambda s: slide_rq2(s, figs),                     # 17
        lambda s: slide_rq3_counter(s, figs),             # 18
        lambda s: slide_rq4_ablation(s, figs),            # 19
        slide_summary,                                    # 20
    ]
    assert len(builders) == EXPECTED_COUNT

    for slide, builder in zip(prs.slides, builders):
        clear_slide(slide)
        builder(slide)
    prs.save(OUT)

    with zipfile.ZipFile(OUT) as zf:
        bad = zf.testzip()
        if bad:
            raise RuntimeError(f"Corrupt pptx member: {bad}")
    check = Presentation(OUT)
    print(f"generated={OUT}")
    print(f"slides={len(check.slides)}")
    print(f"figures={len(figs)}")
    for i, slide in enumerate(check.slides, 1):
        title = ""
        pics = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pics += 1
            if not title and hasattr(shape, "text") and shape.text.strip():
                title = " ".join(shape.text.split())[:64]
        print(f"{i:02d} pics={pics:2d} | {title}")


if __name__ == "__main__":
    main()
