#!/usr/bin/env python3
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation


EXPECTED_TITLES = [
    "穷究于理 成就于工",
    "答辩主线:五个问题串起证据链",
    "研究对象:支持多写的 OLTP 数据库一体机",
    "瓶颈:远程所有权转移如何进入事务关键路径",
    "三类现有方案都不足以解决问题",
    "概念解耦:逻辑/物理位置 + 关键性能代理指标",
    "观察:事务采样 + 元组级亲和图建模",
    "从亲和图到 ParMETIS 划分目标",
    "决策:分布式图划分 + ParMETIS adaptive",
    "发布:快照式 AssignmentTable + epoch 版本号",
    "执行:元组迁移协议(BLink 切换 = 线性化点)",
    "正确性:5 条不变量",
    "性能优化:按 (tableId, srcPage, dstNode) 合并批量迁移",
    "实验设置与四个研究问题",
    "RQ1:在线亲和提升吞吐与尾延迟",
    "RQ1:主对照详细数据(同条件 10 min)",
    "RQ2:6 小时长时间稳定性",
    "RQ3:完成更多迁移 ≠ 更好局部性",
    "RQ4:消融实验 — 批量 / 路由 / 衰减 / 周期",
    "局限、工作总结与未来展望",
]

REQUIRED_PHRASES = [
    "AssignmentTable",
    "24.69",          # Schism 端到端 P99 高 24.69%
    "99.070",         # 在线方案迁移完成率
    "99.597",         # Schism 迁移完成率(对比反直觉)
    "19.660",         # 远程访问比例
    "28912.06",       # 主实验在线方案 TPS
    "34 356.98",      # 6h 平均吞吐
    "50 ms",          # 聚合周期要在主线/机制中出现
    "10 s",           # 重分区周期/周期消融
    "五个问题",        # 故事线必须明确
    "证据链",          # 叙事必须是证据链
    "设计约束",        # 现有方案页需要推出方案约束
    "低侵入",          # 发布路径逻辑
    "故事回扣",        # 总结页必须回扣主线
    "EdgeCut",         # 图划分目标必须解释
    "k(k−1)/2",        # 普通图建模取舍
    "maxChangedVerticesRatio",  # 控制过大迁移
    "summary.txt",     # 实验数据来源口径
    "SmallBank-Aff",   # 实验负载边界
    "waitLockSuccess", # 后续热点页面协调方向
    "图5-9",          # 衰减因子消融必须出现
    "WAL",            # 实验口径声明
]


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return " ".join(shape.text.split())


def first_text(slide) -> str:
    for shape in slide.shapes:
        text = shape_text(shape)
        if text:
            return text
    return ""


def all_text(prs: Presentation) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            t = shape_text(shape)
            if t:
                parts.append(t)
    return "\n".join(parts)


def verify(path: Path) -> None:
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip()
        assert bad is None, f"Corrupt pptx member: {bad}"

    prs = Presentation(path)
    assert len(prs.slides) == 20, f"Expected 20 slides, got {len(prs.slides)}"

    titles = [first_text(slide) for slide in prs.slides]
    for actual, expected in zip(titles, EXPECTED_TITLES):
        assert actual == expected, f"Title mismatch:\n  got:      {actual!r}\n  expected: {expected!r}"

    width, height = prs.slide_width, prs.slide_height
    picture_count = 0
    for idx, slide in enumerate(prs.slides, 1):
        visible_text_count = 0
        for shape in slide.shapes:
            assert shape.left >= -1000, f"Slide {idx}: shape left out of bounds"
            assert shape.top >= -1000, f"Slide {idx}: shape top out of bounds"
            assert shape.left + shape.width <= width + 1000, \
                f"Slide {idx}: shape exceeds right edge"
            assert shape.top + shape.height <= height + 1000, \
                f"Slide {idx}: shape exceeds bottom edge"
            text = shape_text(shape)
            if text:
                visible_text_count += 1
            for bad in ("请输入", "XX研究", "XXX"):
                assert bad not in text, f"Slide {idx}: template placeholder {bad!r} remains"
            if shape.shape_type == 13:
                picture_count += 1
        assert visible_text_count >= 3, f"Slide {idx}: too little visible text ({visible_text_count})"
    assert picture_count >= 15, \
        f"Expected image-led deck with at least 15 pictures, got {picture_count}"

    haystack = all_text(prs)
    for phrase in REQUIRED_PHRASES:
        assert phrase in haystack, f"Required phrase missing from deck: {phrase!r}"


def main() -> None:
    path = Path(sys.argv[1]) if len(sys.argv) > 1 \
        else Path("ppt/支持多写OLTP数据库一体机页面亲和性答辩.pptx")
    verify(path)
    print(f"verify_ok {path}")


if __name__ == "__main__":
    main()
